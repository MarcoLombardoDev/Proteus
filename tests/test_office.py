#!/usr/bin/env python
# Proteus — Rebranding Tool
# Copyright (C) 2026 Marco Lombardo
#
# SPDX-License-Identifier: AGPL-3.0-or-later
# Distributed WITHOUT ANY WARRANTY; see LICENSE for the full terms.
# A commercial licence, without the AGPL's obligations, is available for use
# in proprietary or closed-source products — see COMMERCIAL-LICENSE.md.

"""
Tests for finding and replacing logos inside Office documents.

The documents here are built with the official python-docx / python-pptx /
openpyxl libraries and re-opened with them afterwards. Hand-rolled ZIPs would
prove nothing: the question is whether Word and PowerPoint still accept the
file once Proteus has rewritten it, and the closest available answer is
whether their reference implementations do.

Those libraries are test-only. Proteus itself reads and writes these packages
with `zipfile` alone.
"""

from __future__ import annotations

import os
import sys
import zipfile

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import core  # noqa: E402
import office  # noqa: E402

pytest.importorskip("PIL", reason="Pillow is needed to build test images")
from PIL import Image, ImageDraw  # noqa: E402

docx = pytest.importorskip("docx", reason="python-docx is needed to build documents")
pptx = pytest.importorskip("pptx", reason="python-pptx is needed to build decks")
openpyxl = pytest.importorskip("openpyxl", reason="openpyxl is needed for workbooks")


# ---------------------------------------------------------------------------
# Fixtures: real documents containing a real logo
# ---------------------------------------------------------------------------

def mark(path, size=(240, 80), colour=(196, 62, 58)):
    """A logo-like shape with enough structure to hash."""
    path = str(path)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    image = Image.new("RGB", size, (255, 255, 255))
    draw = ImageDraw.Draw(image)
    w, h = size
    draw.ellipse((w * 0.04, h * 0.15, w * 0.32, h * 0.85), fill=colour)
    draw.rectangle((w * 0.38, h * 0.30, w * 0.94, h * 0.48), fill=colour)
    draw.rectangle((w * 0.38, h * 0.56, w * 0.72, h * 0.72), fill=(90, 90, 90))
    image.save(path)
    return path


def bands(path):
    """
    A genuinely different picture: horizontal bands, not a mark.

    Recolouring the mark would not do — see
    `test_recolouring_still_matches`: the hash reads luminance gradients, so
    the same silhouette in another colour still matches.
    """
    path = str(path)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    image = Image.new("RGB", (240, 80), (255, 255, 255))
    draw = ImageDraw.Draw(image)
    for index in range(6):
        top = 80 * index / 6
        shade = 20 + index * 35
        draw.rectangle((0, top, 240, top + 80 / 6), fill=(shade, shade, shade))
    image.save(path)
    return path


def make_docx(path, logo, in_header=True):
    document = docx.Document()
    document.add_heading("Company report", 0)
    document.add_picture(logo)
    if in_header:
        run = document.sections[0].header.paragraphs[0].add_run()
        run.add_picture(logo)
    document.add_paragraph("Body text that must survive.")
    os.makedirs(os.path.dirname(str(path)), exist_ok=True)
    document.save(str(path))
    return str(path)


def make_pptx(path, logo):
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.add_picture(logo, Inches(1), Inches(1))
    os.makedirs(os.path.dirname(str(path)), exist_ok=True)
    presentation.save(str(path))
    return str(path)


def make_xlsx(path, logo):
    from openpyxl.drawing.image import Image as XLImage

    workbook = openpyxl.Workbook()
    workbook.active["A1"] = "Budget"
    workbook.active.add_image(XLImage(logo), "C3")
    os.makedirs(os.path.dirname(str(path)), exist_ok=True)
    workbook.save(str(path))
    return str(path)


@pytest.fixture
def documents(tmp_path):
    """A folder with one of each format, all carrying the same old logo."""
    old = mark(tmp_path / "brand" / "old_logo.png")
    new = mark(tmp_path / "new" / "new_logo.png", colour=(20, 90, 170))
    folder = tmp_path / "share"
    return {
        "old": old,
        "new": new,
        "folder": str(folder),
        "docx": make_docx(folder / "reports" / "report.docx", old),
        "pptx": make_pptx(folder / "decks" / "deck.pptx", old),
        "xlsx": make_xlsx(folder / "finance" / "budget.xlsx", old),
    }


# ---------------------------------------------------------------------------
# Detection and listing
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("name,expected", [
    ("report.docx", True), ("macro.docm", True), ("template.dotx", True),
    ("deck.pptx", True), ("show.ppsx", True),
    ("budget.xlsx", True), ("sheet.xlsm", True),
    ("logo.png", False), ("notes.txt", False),
    ("legacy.doc", False), ("legacy.ppt", False), ("legacy.xls", False),
])
def test_office_documents_are_recognised_by_extension(name, expected):
    """The pre-2007 binary formats are OLE files, not packages: out of scope."""
    assert office.is_office_document(name) is expected


def test_pictures_are_found_in_every_format(documents):
    for kind, prefix in (("docx", "word/media/"), ("pptx", "ppt/media/"),
                         ("xlsx", "xl/media/")):
        images, _ = office.list_images(documents[kind])
        assert images, f"no picture found in the {kind}"
        assert all(image.entry.startswith(prefix) for image in images)
        assert all(image.size > 0 for image in images)


def test_a_picture_used_twice_is_stored_once(documents):
    """
    The logo sits in the body *and* the header of the report, yet Office keeps
    a single copy — so one replacement fixes both occurrences.
    """
    images, _ = office.list_images(documents["docx"])
    assert len(images) == 1

    with zipfile.ZipFile(documents["docx"]) as package:
        references = sum(package.read(name).count(b"image1.png")
                         for name in package.namelist() if name.endswith(".rels"))
    assert references >= 2


def test_listing_a_non_package_returns_nothing(tmp_path):
    """A corrupt or password-protected file is skipped, not fatal."""
    broken = tmp_path / "broken.docx"
    broken.write_bytes(b"this is not a zip")
    assert office.list_images(str(broken))[0] == []
    assert office.list_images(str(tmp_path / "missing.docx"))[0] == []


def test_metafiles_are_ignored(tmp_path, documents):
    """EMF/WMF cannot be read by Pillow, so they must not be offered."""
    target = str(tmp_path / "with_emf.docx")
    with zipfile.ZipFile(documents["docx"]) as src, \
            zipfile.ZipFile(target, "w") as dst:
        for info in src.infolist():
            dst.writestr(info.filename, src.read(info.filename))
        dst.writestr("word/media/image9.emf", b"fake metafile")

    entries = [image.entry for image in office.list_images(target)[0]]
    assert "word/media/image9.emf" not in entries


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------

def test_embedded_pictures_can_be_extracted_and_hashed(documents):
    """Content search has to reach inside documents, and this is how."""
    old_hash = core.perceptual_hash(documents["old"])

    for kind in ("docx", "pptx", "xlsx"):
        image = office.list_images(documents[kind])[0][0]
        temp = office.extract_to_temp(image.document, image.entry)
        assert temp and os.path.exists(temp)
        try:
            assert core.hash_similarity(core.perceptual_hash(temp), old_hash) == 1.0
        finally:
            os.remove(temp)


def test_extracting_a_missing_entry_is_not_fatal(documents):
    assert office.extract(documents["docx"], "word/media/nope.png") is None
    assert office.extract_to_temp(documents["docx"], "word/media/nope.png") is None


def test_keys_round_trip(documents):
    image = office.list_images(documents["docx"])[0][0]
    assert office.is_embedded_key(image.key)
    assert office.split_key(image.key) == (image.document, image.entry)
    assert office.split_key("/plain/file.png") is None
    assert image.name.endswith("image1.png")


# ---------------------------------------------------------------------------
# Replacement
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("kind", ["docx", "pptx", "xlsx"])
def test_replacement_changes_the_picture_and_keeps_the_document_valid(documents, kind):
    document = documents[kind]
    entry = office.list_images(document)[0][0].entry
    new_bytes = open(documents["new"], "rb").read()

    office.write_replacements(document, {entry: new_bytes})

    # The picture really changed...
    temp = office.extract_to_temp(document, entry)
    try:
        assert Image.open(temp).convert("RGB").getpixel((60, 40)) == (20, 90, 170)
    finally:
        os.remove(temp)

    # ...and the official library still reads the document.
    if kind == "docx":
        reopened = docx.Document(document)
        assert reopened.paragraphs[0].text == "Company report"
        assert len(reopened.inline_shapes) == 1
        assert reopened.sections[0].header.paragraphs
    elif kind == "pptx":
        reopened = pptx.Presentation(document)
        assert len(reopened.slides) == 1
        assert any(shape.shape_type == 13 for shape in reopened.slides[0].shapes)
    else:
        reopened = openpyxl.load_workbook(document)
        assert reopened.active["A1"].value == "Budget"
        assert len(reopened.active._images) == 1


def test_untouched_entries_are_copied_byte_for_byte(documents):
    document = documents["docx"]
    entry = office.list_images(document)[0][0].entry

    with zipfile.ZipFile(document) as package:
        before = {name: package.read(name) for name in package.namelist()}

    office.write_replacements(document, {entry: b"\x89PNG\r\n\x1a\nreplaced"})

    with zipfile.ZipFile(document) as package:
        after = {name: package.read(name) for name in package.namelist()}

    assert set(before) == set(after), "no entry may appear or vanish"
    for name in before:
        if name != entry:
            assert before[name] == after[name], f"{name} was altered"


def test_entry_order_is_preserved(documents):
    """Some readers care where [Content_Types].xml sits in the package."""
    document = documents["docx"]
    with zipfile.ZipFile(document) as package:
        before = package.namelist()

    office.write_replacements(document, {before[-1]: b"x"}
                              if before[-1].endswith(".png") else {})
    with zipfile.ZipFile(document) as package:
        assert package.namelist() == before


def test_a_failed_rewrite_leaves_the_document_intact(documents, monkeypatch):
    """The package is rebuilt, so a failure must not destroy the original."""
    document = documents["docx"]
    entry = office.list_images(document)[0][0].entry
    original = open(document, "rb").read()

    real_writestr = zipfile.ZipFile.writestr

    def exploding(self, zinfo, data, *args, **kwargs):
        if getattr(zinfo, "filename", zinfo) == entry:
            raise OSError("disk full")
        return real_writestr(self, zinfo, data, *args, **kwargs)

    monkeypatch.setattr(zipfile.ZipFile, "writestr", exploding)

    with pytest.raises(OSError):
        office.write_replacements(document, {entry: b"new"})

    assert open(document, "rb").read() == original
    leftovers = [f for f in os.listdir(os.path.dirname(document))
                 if f.startswith(".proteus_")]
    assert leftovers == []


# ---------------------------------------------------------------------------
# Aspect ratio: the silent way to ruin a document
# ---------------------------------------------------------------------------

def test_aspect_mismatch_flags_a_distorting_replacement():
    assert office.aspect_mismatch((240, 80), (200, 200)) is True   # 3:1 -> 1:1
    assert office.aspect_mismatch((240, 80), (480, 160)) is False  # same ratio
    assert office.aspect_mismatch((240, 80), (243, 80)) is False   # within tolerance
    assert office.aspect_mismatch((240, 80), (300, 80)) is True    # 3:1 -> 3.75:1


def test_aspect_mismatch_is_silent_when_a_size_is_unknown():
    assert office.aspect_mismatch(None, (200, 200)) is False
    assert office.aspect_mismatch((240, 80), None) is False
    assert office.aspect_ratio((0, 0)) is None


def test_the_frame_does_not_follow_the_picture(documents):
    """
    Why the check above exists: the shape geometry lives in the document XML,
    not in the image, so a square logo dropped into a 3:1 frame is stretched
    rather than resized. Nothing errors; the document just looks wrong.
    """
    document = documents["docx"]
    before = docx.Document(document).inline_shapes[0]
    ratio_before = before.width / before.height

    square = mark(os.path.join(os.path.dirname(documents["new"]), "square.png"),
                  size=(200, 200), colour=(20, 90, 170))
    entry = office.list_images(document)[0][0].entry
    office.write_replacements(document, {entry: open(square, "rb").read()})

    after = docx.Document(document).inline_shapes[0]
    assert after.width / after.height == pytest.approx(ratio_before), (
        "the frame keeps its old proportions, which is exactly why "
        "aspect_mismatch has to warn before this happens"
    )


# ---------------------------------------------------------------------------
# Integration with the scan and the replacement pipeline
# ---------------------------------------------------------------------------

def test_scanning_a_folder_finds_pictures_inside_documents(documents):
    found = core.scan_office_documents(documents["folder"])
    assert len(found) == 3
    assert all(info.embedded for info in found)
    assert all(info.dim == (240, 80) for info in found)
    assert sorted(os.path.basename(info.container) for info in found) == [
        "budget.xlsx", "deck.pptx", "report.docx"]


def test_scanning_documents_by_content_filters_out_other_pictures(documents, tmp_path):
    """A document holding an unrelated picture must not be picked up."""
    other = bands(tmp_path / "other" / "photo.png")
    unrelated = make_docx(tmp_path / "share" / "reports" / "unrelated.docx",
                          other, in_header=False)
    assert os.path.exists(unrelated)

    everything = core.scan_office_documents(documents["folder"])
    assert len(everything) == 4

    matching = core.scan_office_documents(documents["folder"],
                                          references=[documents["old"]])
    assert len(matching) == 3
    assert all(info.similarity == 1.0 for info in matching)
    assert not any("unrelated" in info.container for info in matching)


def test_replacement_through_the_normal_pipeline(documents):
    """Embedded pictures flow through build_matches and replace_all unchanged."""
    targets = core.scan_office_documents(documents["folder"],
                                         references=[documents["old"]])
    sources = [core.FileInfo.from_path(documents["new"])]

    matches = core.build_matches(targets, sources)
    assert all(match.source is not None for match in matches)

    report = core.replace_all(matches, backup=True)
    assert report.ok == 3 and report.errors == 0

    for kind in ("docx", "pptx", "xlsx"):
        entry = office.list_images(documents[kind])[0][0].entry
        temp = office.extract_to_temp(documents[kind], entry)
        try:
            assert Image.open(temp).convert("RGB").getpixel((60, 40)) == (20, 90, 170)
        finally:
            os.remove(temp)
        assert os.path.exists(documents[kind] + ".bak")


def test_one_document_is_backed_up_once_however_many_pictures_it_holds(tmp_path):
    """
    Replacing three logos in one report must not produce three backups of
    successive states — that would leave no clean copy of the original.
    """
    old = mark(tmp_path / "brand" / "old.png")
    new = mark(tmp_path / "new" / "new.png", colour=(20, 90, 170))

    # Office stores pictures by content, so three *identical* images would be
    # deduplicated into a single media entry. They have to differ.
    document = docx.Document()
    for index, size in enumerate([(240, 80), (200, 120), (160, 160)]):
        document.add_picture(mark(tmp_path / "brand" / f"copy{index}.png", size=size))
    path = str(tmp_path / "share" / "many.docx")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    document.save(path)

    targets = core.scan_office_documents(str(tmp_path / "share"))
    assert len(targets) == 3, "three distinct pictures inside one document"

    sources = [core.FileInfo.from_path(new)]
    report = core.replace_all(core.build_matches(targets, sources), backup=True)

    assert report.total == 1, "one document, one outcome"
    assert report.ok == 1
    backups = [f for f in os.listdir(os.path.dirname(path)) if f.endswith(".bak")]
    assert len(backups) == 1, f"expected a single backup, found {backups}"

    # And that single backup is the untouched original.
    assert core.perceptual_hash(old) is not None
    restored = office.extract_to_temp(os.path.join(os.path.dirname(path), backups[0]),
                                      "word/media/image1.png")
    try:
        assert Image.open(restored).convert("RGB").getpixel((60, 40)) != (20, 90, 170)
    finally:
        os.remove(restored)


def test_dry_run_does_not_touch_documents(documents):
    before = open(documents["docx"], "rb").read()
    targets = core.scan_office_documents(documents["folder"])
    sources = [core.FileInfo.from_path(documents["new"])]

    report = core.replace_all(core.build_matches(targets, sources),
                              backup=True, dry_run=True)
    assert report.ok == 3
    assert open(documents["docx"], "rb").read() == before
    assert not os.path.exists(documents["docx"] + ".bak")


def test_restore_brings_a_document_back(documents):
    targets = core.scan_office_documents(documents["folder"])
    sources = [core.FileInfo.from_path(documents["new"])]
    core.replace_all(core.build_matches(targets, sources), backup=True)

    report = core.restore_backups(documents["folder"])
    assert report.ok == 3

    entry = office.list_images(documents["docx"])[0][0].entry
    temp = office.extract_to_temp(documents["docx"], entry)
    try:
        assert Image.open(temp).convert("RGB").getpixel((60, 40)) == (196, 62, 58)
    finally:
        os.remove(temp)


def test_a_missing_document_is_reported_not_raised(tmp_path):
    outcome = core.replace_in_document(str(tmp_path / "gone.docx"),
                                       {"word/media/image1.png": "whatever"})
    assert outcome.status == "error"


# ---------------------------------------------------------------------------
# Reporting what cannot be done
# ---------------------------------------------------------------------------

def test_a_pasted_logo_is_reported_not_ignored(tmp_path):
    """
    Regression, and the commonest case of all.

    Office stores a *pasted* logo as an EMF or WMF metafile. Pillow cannot
    rasterise those, so they can be neither compared nor replaced — but they are
    how most logos get into a corporate Word document. Returning nothing at all
    left the commonest case silently unhandled, which is exactly what the
    "nothing is skipped in silence" rule forbids.
    """
    import zipfile

    document = str(tmp_path / "pasted.docx")
    with zipfile.ZipFile(document, "w") as package:
        package.writestr("[Content_Types].xml", "<Types/>")
        package.writestr("word/document.xml", "<w:document/>")
        package.writestr("word/media/image1.emf", b"\x01\x00\x00\x00fake-emf")

    images, problems = office.list_images(document)

    assert images == [], "an EMF cannot be compared or replaced"
    assert len(problems) == 1
    assert "image1.emf" in problems[0].reason
    assert problems[0].hint, "a finding without a remedy is only half reported"


def test_a_password_protected_document_says_so(tmp_path):
    """
    A protected .docx is an OLE compound file, not a ZIP, so it merely looks
    corrupt. Saying "damaged" would send the user hunting the wrong problem.
    """
    document = str(tmp_path / "locked.docx")
    with open(document, "wb") as handle:
        handle.write(office.OLE_MAGIC + b"rest of an OLE container")

    images, problems = office.list_images(document)

    assert images == []
    assert len(problems) == 1
    assert "password" in problems[0].reason.lower()


def test_a_damaged_document_is_reported(tmp_path):
    document = str(tmp_path / "broken.docx")
    with open(document, "wb") as handle:
        handle.write(b"PK\x03\x04 truncated")

    images, problems = office.list_images(document)
    assert images == []
    assert problems and problems[0].hint


def test_a_clean_document_reports_nothing(documents):
    """The bar has to stay silent when there is nothing to say."""
    # The fixture also carries plain image paths and folders; only the packages
    # are relevant here.
    for kind in ("docx", "pptx", "xlsx"):
        images, problems = office.list_images(documents[kind])
        assert images, kind
        assert problems == [], kind


def test_the_scan_forwards_office_findings(tmp_path):
    """End to end: the finding reaches the caller, not just the module."""
    import zipfile

    import core

    folder = tmp_path / "share"
    os.makedirs(folder, exist_ok=True)
    with zipfile.ZipFile(str(folder / "pasted.docx"), "w") as package:
        package.writestr("[Content_Types].xml", "<Types/>")
        package.writestr("word/media/logo.wmf", b"fake-wmf")

    seen = []
    found = core.scan_office_documents(str(folder), on_problem=seen.append)

    assert found == []
    assert len(seen) == 1 and "logo.wmf" in seen[0].reason
